from abc import ABC, abstractmethod
import re
from typing import List, Union, Dict
from enum import Enum

class TypeParser(Enum):
    PDF = "application/pdf"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    TXT = "text/plain"

class DocumentParser(ABC):
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.text = ""
        self.chunks = []

    @abstractmethod
    def get_text(self) -> str:
        """Extracts all text from the document."""
        pass

    def split_text(self, text: str, chunk_size: int = 200) -> list:
        """
        Splits the text into chunks of nearly `chunk_size` words.
        Returns a list of text chunks.
        """
        words = re.findall(r'\S+', text)
        self.chunks = []
        for i in range(0, len(words), chunk_size):
            chunk = ' '.join(words[i:i+chunk_size])
            self.chunks.append(chunk)
        return self.chunks

class PdfParser(DocumentParser):
    def get_text(self) -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(self.file_path)
        self.text = ""
        for page in reader.pages:
            self.text += page.extract_text() or ""
        return self.text

class PptxParser(DocumentParser):
    def get_text(self) -> str:
        from pptx import Presentation
        prs = Presentation(self.file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text

class DocxParser(DocumentParser):
    def get_text(self) -> str:
        from docx import Document
        doc = Document(self.file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + " "
        return text

type2Class = {
    TypeParser.PDF: PdfParser,
    TypeParser.PPTX: PptxParser,
    TypeParser.DOCX: DocxParser
}

def pipeline_parser(file_obj) -> List[Dict[str, Union[str, TypeParser]]]:
    """
    Returns a list of chunk dicts for a Chainlit file or AskFileResponse.
    """
    file_type = None
    file_path = None
    file_id = None
    file_name = None

    # Chainlit File: has .mime, .path, .id, .name
    # AskFileResponse: has .type, .path, .id, .name
    if file_obj["type"]:
        # Chainlit File
        mime = file_obj["type"]
        if mime:
            if "pdf" in mime:
                file_type = TypeParser.PDF
            elif "presentation" in mime or "pptx" in mime:
                file_type = TypeParser.PPTX
            elif "word" in mime or "docx" in mime:
                file_type = TypeParser.DOCX
        file_path = file_obj["path"]
        file_id = file_obj["id"]
        file_name = file_obj["name"]

    if file_type in type2Class and file_path:
        parsing_class: DocumentParser = type2Class[file_type]
        parser = parsing_class(file_path)
        parser.get_text()
        parser.split_text(parser.text)
        return [{
            "id": file_id,
            "file_path": file_path,
            "file_type": file_type,
            "file_name": file_name,
            "chunk": chunk,
        } for chunk in parser.chunks]
    else:
        raise ValueError(f"Unsupported or missing file type: {file_type}, path: {file_path}")